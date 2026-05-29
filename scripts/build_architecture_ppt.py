#!/usr/bin/env python3

from __future__ import annotations

import sqlite3
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ARTIFACT_DIR = ROOT / "artifacts" / "architecture"
SCREENSHOT_DIR = ARTIFACT_DIR / "screenshots"
DB_PATH = ROOT / "data" / "app.db"
OUTPUT_PATH = ARTIFACT_DIR / "RRCP_Architecture_Walkthrough.pptx"


NAVY = RGBColor(30, 42, 74)
BLUE = RGBColor(49, 94, 251)
SOFT_BLUE = RGBColor(236, 242, 255)
TEXT = RGBColor(39, 50, 74)
MUTED = RGBColor(95, 112, 148)
GREEN = RGBColor(42, 155, 96)
AMBER = RGBColor(210, 140, 24)
RED = RGBColor(194, 63, 63)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(210, 220, 238)
PALE = RGBColor(247, 249, 253)


def counts_snapshot() -> dict[str, int]:
    snapshot: dict[str, int] = {}
    if not DB_PATH.exists():
        return snapshot

    with sqlite3.connect(DB_PATH) as conn:
        for table in [
            "users",
            "doctor_master",
            "service_prices",
            "discount_rules",
            "software_requirements",
            "transactions",
            "engine_runs",
            "engine_results",
            "payments",
            "approval_requests",
            "locked_periods",
            "contact_messages",
        ]:
            snapshot[table] = conn.execute(f"SELECT COUNT(*) FROM {table}").fetchone()[0]
    return snapshot


def add_full_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide, title: str, subtitle: str):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.65),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.16), Inches(7.4), Inches(0.25))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = WHITE

    sub_box = slide.shapes.add_textbox(Inches(8.0), Inches(0.17), Inches(4.8), Inches(0.24))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Aptos"
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(216, 226, 250)


def add_page_title(slide, title: str, subtitle: str):
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.9), Inches(6.5), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = NAVY

    sub_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.38), Inches(6.6), Inches(0.45))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Aptos"
    run.font.size = Pt(12.5)
    run.font.color.rgb = MUTED


def add_body_card(slide, x, y, w, h, title: str, bullets: list[str], accent: RGBColor = BLUE):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE
    card.line.width = Pt(1)

    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.color.rgb = accent

    title_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.14), w - Inches(0.3), Inches(0.32))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos"
    run.font.bold = True
    run.font.size = Pt(15)
    run.font.color.rgb = TEXT

    body = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.48), w - Inches(0.32), h - Inches(0.56))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_top = 0
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT
        p.space_after = Pt(7)
        p.bullet = True


def fit_picture(slide, path: Path, x, y, w, h):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = SOFT_BLUE
    box.line.color.rgb = LINE
    box.line.width = Pt(1)

    with Image.open(path) as image:
        img_w, img_h = image.size

    target_ratio = w / h
    image_ratio = img_w / img_h
    if image_ratio > target_ratio:
        pic_w = w - Inches(0.18)
        pic_h = pic_w / image_ratio
    else:
        pic_h = h - Inches(0.18)
        pic_w = pic_h * image_ratio

    pic_x = x + (w - pic_w) / 2
    pic_y = y + (h - pic_h) / 2
    slide.shapes.add_picture(str(path), pic_x, pic_y, width=pic_w, height=pic_h)


def add_metric_chip(slide, x, y, w, label: str, value: str, accent: RGBColor = BLUE):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.72))
    chip.fill.solid()
    chip.fill.fore_color.rgb = WHITE
    chip.line.color.rgb = LINE
    chip.line.width = Pt(1)

    accent_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.1), y + Inches(0.13), Inches(0.18), Inches(0.18))
    accent_box.fill.solid()
    accent_box.fill.fore_color.rgb = accent
    accent_box.line.color.rgb = accent

    label_box = slide.shapes.add_textbox(x + Inches(0.38), y + Inches(0.08), w - Inches(0.45), Inches(0.2))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.name = "Aptos"
    run.font.size = Pt(9.5)
    run.font.color.rgb = MUTED

    value_box = slide.shapes.add_textbox(x + Inches(0.38), y + Inches(0.28), w - Inches(0.45), Inches(0.24))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = value
    run.font.name = "Aptos"
    run.font.bold = True
    run.font.size = Pt(14)
    run.font.color.rgb = TEXT


def add_small_box(slide, x, y, w, h, title: str, body: str, fill=WHITE, accent=BLUE):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = LINE
    box.line.width = Pt(1)

    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.12), y + Inches(0.14), Inches(0.16), Inches(0.16))
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent
    dot.line.color.rgb = accent

    title_box = slide.shapes.add_textbox(x + Inches(0.34), y + Inches(0.08), w - Inches(0.44), Inches(0.22))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos"
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.color.rgb = TEXT

    body_box = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.34), w - Inches(0.24), h - Inches(0.42))
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.name = "Aptos"
    run.font.size = Pt(9.5)
    run.font.color.rgb = MUTED


def add_arrow(slide, x, y, w, h, text: str):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = SOFT_BLUE
    arrow.line.color.rgb = BLUE
    arrow.line.width = Pt(1)

    box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.08), w - Inches(0.25), h - Inches(0.12))
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.color.rgb = NAVY


def add_table_like(slide, x, y, w, headers: list[str], rows: list[list[str]], col_widths: list[float]):
    row_h = 0.36
    total_h = Inches(row_h * (len(rows) + 1) + 0.04)
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, total_h)
    outer.fill.solid()
    outer.fill.fore_color.rgb = WHITE
    outer.line.color.rgb = LINE
    outer.line.width = Pt(1)

    cur_x = x
    for idx, (header, ratio) in enumerate(zip(headers, col_widths)):
        col_w = Inches(ratio)
        fill_color = SOFT_BLUE
        cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cur_x, y, col_w, Inches(row_h))
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
        cell.line.color.rgb = LINE
        cell.line.width = Pt(0.5)
        box = slide.shapes.add_textbox(cur_x + Inches(0.06), y + Inches(0.06), col_w - Inches(0.12), Inches(0.2))
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = header
        run.font.name = "Aptos"
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = NAVY
        cur_x += col_w

    for row_idx, row in enumerate(rows, start=1):
        cur_x = x
        row_y = y + Inches(row_h * row_idx)
        for value, ratio in zip(row, col_widths):
            col_w = Inches(ratio)
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cur_x, row_y, col_w, Inches(row_h))
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 else PALE
            cell.line.color.rgb = LINE
            cell.line.width = Pt(0.5)
            box = slide.shapes.add_textbox(cur_x + Inches(0.06), row_y + Inches(0.05), col_w - Inches(0.12), Inches(0.22))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = value
            run.font.name = "Aptos"
            run.font.size = Pt(9.5)
            run.font.color.rgb = TEXT
            cur_x += col_w


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    counts = counts_snapshot()

    overview_png = SCREENSHOT_DIR / "overview.png"
    setup_png = SCREENSHOT_DIR / "setup-center.png"
    intake_png = SCREENSHOT_DIR / "monthly-intake.png"
    calc_png = SCREENSHOT_DIR / "calculation-review.png"
    payout_png = SCREENSHOT_DIR / "payout-center.png"
    reports_png = SCREENSHOT_DIR / "reports.png"
    support_png = SCREENSHOT_DIR / "support.png"

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, PALE)
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.45), Inches(12.43), Inches(6.55))
    hero.fill.solid()
    hero.fill.fore_color.rgb = WHITE
    hero.line.color.rgb = LINE
    hero.line.width = Pt(1)
    add_top_band(slide, "Referral Revenue Calculation Platform", "Architecture walk-through | March 16, 2026")

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.15), Inches(5.2), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Full workflow architecture"
    run.font.name = "Aptos Display"
    run.font.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = NAVY
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Frontend, backend, database, approvals, page-level functioning, and end-to-end monthly data movement."
    run.font.name = "Aptos"
    run.font.size = Pt(13)
    run.font.color.rgb = MUTED

    bullets = slide.shapes.add_textbox(Inches(0.92), Inches(2.1), Inches(4.9), Inches(2.5))
    tf = bullets.text_frame
    tf.word_wrap = True
    for idx, text in enumerate(
        [
            "Frontend: index.php + public/app.js + public/styles.css",
            "Backend: api/index.php routing + api/lib.php schema, upload parsing, engine, reporting",
            "Database: SQLite file at data/app.db",
            "Current demo dataset: "
            f"{counts.get('users', 0)} users, {counts.get('doctor_master', 0)} doctors, "
            f"{counts.get('transactions', 0)} transactions, {counts.get('approval_requests', 0)} approvals",
        ]
    ):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.bullet = True
        p.font.name = "Aptos"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)

    add_metric_chip(slide, Inches(0.92), Inches(4.95), Inches(1.95), "Workflow pages", "7 pages", BLUE)
    add_metric_chip(slide, Inches(3.02), Inches(4.95), Inches(1.95), "Roles", "4 roles", GREEN)
    add_metric_chip(slide, Inches(5.12), Inches(4.95), Inches(1.95), "Approval types", "5 live", AMBER)
    fit_picture(slide, overview_png, Inches(7.0), Inches(1.15), Inches(5.2), Inches(4.95))

    footer = slide.shapes.add_textbox(Inches(0.92), Inches(6.2), Inches(5.7), Inches(0.35))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Purpose: explain how the tool starts, where each upload lands, how calculation runs, how approvals are applied, and how the month closes."
    run.font.name = "Aptos"
    run.font.size = Pt(10.5)
    run.font.color.rgb = MUTED

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide)
    add_top_band(slide, "System Architecture", "Component layout and runtime responsibilities")
    add_page_title(slide, "How the application is structured", "The product is a single-page frontend on top of PHP API endpoints and a local SQLite operational database.")
    add_small_box(slide, Inches(0.7), Inches(2.0), Inches(2.3), Inches(1.55), "Client / UI layer", "index.php boots the app shell.\npublic/app.js renders workflow pages and calls the API.\npublic/styles.css defines the compact responsive UI.", fill=WHITE, accent=BLUE)
    add_small_box(slide, Inches(3.4), Inches(2.0), Inches(2.45), Inches(1.55), "API / routing layer", "api/index.php authenticates requests, enforces role access, handles uploads, exports, approvals, and payment actions.", fill=WHITE, accent=GREEN)
    add_small_box(slide, Inches(6.25), Inches(2.0), Inches(2.45), Inches(1.55), "Business logic layer", "api/lib.php owns schema creation, seed/demo data, Excel parsing, RRCP engine computation, payment generation, and migrations.", fill=WHITE, accent=AMBER)
    add_small_box(slide, Inches(9.1), Inches(2.0), Inches(3.0), Inches(1.55), "Persistence layer", "data/app.db stores operational tables.\nuploads/ stores incoming workbooks.\ntemplates/ stores sample master formats.", fill=WHITE, accent=RED)
    add_arrow(slide, Inches(2.95), Inches(2.48), Inches(0.38), Inches(0.48), "API")
    add_arrow(slide, Inches(5.82), Inches(2.48), Inches(0.38), Inches(0.48), "Logic")
    add_arrow(slide, Inches(8.68), Inches(2.48), Inches(0.38), Inches(0.48), "DB")

    add_body_card(
        slide,
        Inches(0.7),
        Inches(4.0),
        Inches(5.7),
        Inches(2.55),
        "Supported business inputs",
        [
            "Special Discount Master.xlsx seeds service prices, discount rules, doctor master, incentive groups, and PRO assignments.",
            "Software requirement.xlsx populates the rule checklist / requirement repository used for operating guidance.",
            "Dashboard and incentive transaction workbooks populate normalized monthly intake rows in transactions.",
            "User logins are managed in Setup Center and linked back to doctor_master for doctor-role scope restriction.",
        ],
    )
    add_body_card(
        slide,
        Inches(6.7),
        Inches(4.0),
        Inches(5.95),
        Inches(2.55),
        "Why this layout matters",
        [
            "The frontend remains simple: one JS app, one stylesheet, workflow-first navigation.",
            "The backend remains auditable: approval events and payment edits persist in dedicated tables.",
            "The data model is compact enough for local deployment but structured enough for future migration to MySQL/Postgres if needed.",
            "All monthly operations revolve around a predictable path: setup -> intake -> calculation -> payout -> reporting -> support.",
        ],
        accent=GREEN,
    )

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide)
    add_top_band(slide, "Roles And Workflow", "Who uses what and where their data scope is enforced")
    add_page_title(slide, "Role model and navigation structure", "The app is intentionally workflow-based, so each role sees the same monthly sequence with different permissions.")
    add_table_like(
        slide,
        Inches(0.75),
        Inches(1.95),
        Inches(6.25),
        ["Role", "Main responsibility", "Scope"],
        [
            ["Admin", "Owns setup, approvals, payouts, month lock", "Full access"],
            ["Mapper", "Maintains masters, doctor mapping, monthly intake", "No final admin controls"],
            ["Accountant", "Runs calculation, reviews payouts, exports reports", "No user admin"],
            ["Doctor", "Views own linked data and report output", "Doctor-linked rows only"],
        ],
        [1.2, 3.0, 2.05],
    )
    add_body_card(
        slide,
        Inches(7.3),
        Inches(1.95),
        Inches(5.2),
        Inches(1.95),
        "Workflow navigation",
        [
            "Overview",
            "Setup Center",
            "Monthly Intake",
            "Calculation Review",
            "Payout Center",
            "Reports",
            "Support",
        ],
        accent=BLUE,
    )
    add_body_card(
        slide,
        Inches(7.3),
        Inches(4.1),
        Inches(5.2),
        Inches(2.25),
        "Access control notes",
        [
            "Doctor users must be linked to a doctor_master record; they cannot see unrelated doctors.",
            "Approval actions remain admin-only even when the request was created by mapper or accountant.",
            "Search and filtering are available across the major operational tables to reduce manual scrolling.",
            "Month lock prevents fresh processing on closed periods once finance is complete.",
        ],
        accent=GREEN,
    )

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide)
    add_top_band(slide, "Setup Center", "Master data, team access, doctor governance, and request initiation")
    add_page_title(slide, "Setup Center: where the month is prepared", "This page combines user maintenance, reference uploads, doctor verification, and approval-request creation.")
    add_body_card(
        slide,
        Inches(0.62),
        Inches(1.82),
        Inches(4.35),
        Inches(4.98),
        "What happens here",
        [
            "Create or update users for admin, mapper, accountant, and doctor roles.",
            "Upload Special Discount Master to refresh service_prices, discount_rules, and doctor_master.",
            "Upload Software Requirements to refresh software_requirements.",
            "Verify doctors and maintain fields such as incentive group, incentive cycle, reporting doctor, present PRO, and confirmation status.",
            "Raise approval requests for PRO change, doctor-info change, and doctor addition before admin action.",
        ],
        accent=BLUE,
    )
    fit_picture(slide, setup_png, Inches(5.2), Inches(1.82), Inches(7.55), Inches(4.98))

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide)
    add_top_band(slide, "Monthly Intake", "Normalization, validation, search, and export of transaction rows")
    add_page_title(slide, "Monthly Intake: where raw billing rows enter the system", "Uploaded workbooks are parsed into normalized transaction records and kept searchable before the engine runs.")
    add_body_card(
        slide,
        Inches(0.62),
        Inches(1.82),
        Inches(4.4),
        Inches(5.0),
        "Data path",
        [
            "Input files: dashboard or incentive workbook for the selected month.",
            "Parser maps patient ID, patient name, doctor, PRO, item, price, discount, net, total payment, payment method, and revenue booked center.",
            "Each imported row lands in transactions with raw_json retained for audit and backfill repair.",
            "Users can search across all visible table columns, paginate, and export a filtered CSV for offline review.",
        ],
        accent=GREEN,
    )
    fit_picture(slide, intake_png, Inches(5.15), Inches(1.82), Inches(7.6), Inches(5.0))

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide)
    add_top_band(slide, "Calculation Review", "RRCP engine execution, flags, variance review, and payable computation")
    add_page_title(slide, "Calculation Review: how the monthly engine works", "The engine compares actual discount behavior against uploaded rules and doctor grouping to compute payable incentive.")
    add_body_card(
        slide,
        Inches(0.62),
        Inches(1.82),
        Inches(4.45),
        Inches(5.0),
        "Engine inputs and outputs",
        [
            "Inputs: transactions for the chosen month + doctor_master + discount_rules + service_prices.",
            "Output header: engine_runs stores run summary, counts, and total flags.",
            "Output rows: engine_results stores doctor, group, PRO, modality/item, allowed discount, actual discount, variance, payable amount, and remarks.",
            "Override of incentive amount raises an approval request rather than directly mutating payout value.",
            "Productivity projection gives a quick PRO view of cases and suggested incentive exposure.",
        ],
        accent=AMBER,
    )
    fit_picture(slide, calc_png, Inches(5.2), Inches(1.82), Inches(7.55), Inches(5.0))

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide)
    add_top_band(slide, "Payout Center", "Payment entry creation, disbursal controls, handover tracking, and month lock")
    add_page_title(slide, "Payout Center: from engine output to controlled disbursal", "Payment rows are generated from engine results, then governed through approval, cash-in-hand checks, and timestamps.")
    add_body_card(
        slide,
        Inches(0.62),
        Inches(1.82),
        Inches(4.55),
        Inches(5.0),
        "Payout controls now covered",
        [
            "Creates payment entries from the selected engine run into payments.",
            "Captures adjustment amount, advance payment, return of incentive, PRO cash in hand, manager cash in hand, cashier handover time, and PRO handover time.",
            "Fresh disbursal is blocked when PRO or manager cash in hand is not zero.",
            "Approval of disbursal stays in approval_requests and updates payment approval_status only after admin approval.",
            "Date lock writes into locked_periods once the month is finalized.",
        ],
        accent=RED,
    )
    fit_picture(slide, payout_png, Inches(5.35), Inches(1.82), Inches(7.4), Inches(5.0))

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide)
    add_top_band(slide, "Reports And Support", "Exports, audit outputs, and issue escalation")
    add_page_title(slide, "Reporting and support close the loop", "After payout review, users export audit files or raise structured support requests for correction and follow-up.")
    add_body_card(
        slide,
        Inches(0.62),
        Inches(1.82),
        Inches(4.3),
        Inches(5.0),
        "What these pages do",
        [
            "Reports exports doctor-level or grouped CSV outputs from the latest run in the chosen period.",
            "Doctor users are limited to their own linked individual report.",
            "Support writes structured requests into contact_messages so admin can track mapping issues, payout clarification, and access problems.",
            "These pages are the end-user output surface of the monthly cycle once processing is complete.",
        ],
        accent=BLUE,
    )
    fit_picture(slide, reports_png, Inches(5.2), Inches(1.8), Inches(3.55), Inches(4.9))
    fit_picture(slide, support_png, Inches(8.95), Inches(1.8), Inches(3.55), Inches(4.9))

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide)
    add_top_band(slide, "Approval Architecture", "Requests, admin decisions, and entity updates")
    add_page_title(slide, "Approval process from request to effect", "Approval requests are stored centrally so every exception stays explicit and reviewable.")
    add_table_like(
        slide,
        Inches(0.75),
        Inches(1.9),
        Inches(5.5),
        ["Approval type", "Created from", "Applied to"],
        [
            ["change_of_pro", "Setup Center", "doctor_master.present_pro"],
            ["change_of_doctor_info", "Setup Center", "doctor_master profile fields"],
            ["addition_of_doctor", "Setup Center", "new doctor_master row"],
            ["override_of_incentive_amount", "Calculation Review", "payments.amount after approval"],
            ["approval_of_disbursal", "Payout Center / seeded flow", "payments.approval_status + payment status"],
        ],
        [2.2, 1.45, 1.85],
    )
    add_body_card(
        slide,
        Inches(6.6),
        Inches(1.9),
        Inches(5.85),
        Inches(4.85),
        "Approval lifecycle",
        [
            "1. User raises request with payload_json and requester identity.",
            "2. Request is stored in approval_requests with pending status.",
            "3. Admin reviews readable payload details in the approval table.",
            "4. Approve = backend mutates target entity and stamps approved_by.",
            "5. Reject = request is kept for audit with rejected state, without applying business changes.",
        ],
        accent=AMBER,
    )

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide)
    add_top_band(slide, "Data Model And Movement", "Where data starts, where it lands, and what each table stores")
    add_page_title(slide, "Core tables and how data moves between them", "This is the operational spine of the application. Every page writes to, reads from, or derives from these tables.")
    left_rows = [
        ["users", "Login identities, role, linked doctor"],
        ["doctor_master", "Doctor profile, PRO, group, cycle, verification"],
        ["service_prices", "Uploaded service price list"],
        ["discount_rules", "Discount ceilings per item / group"],
        ["software_requirements", "Operating requirement checklist"],
        ["transactions", "Normalized billing intake rows with raw_json"],
    ]
    right_rows = [
        ["engine_runs", "One summary per calculation run"],
        ["engine_results", "Row-level allowed vs actual discount"],
        ["payments", "Disbursal records and finance adjustments"],
        ["approval_requests", "Exception and approval workflow queue"],
        ["locked_periods", "Month close control"],
        ["contact_messages", "Support inbox items"],
    ]
    add_table_like(slide, Inches(0.72), Inches(1.85), Inches(5.9), ["Table", "Purpose"], left_rows, [1.65, 4.15])
    add_table_like(slide, Inches(6.75), Inches(1.85), Inches(5.85), ["Table", "Purpose"], right_rows, [1.8, 4.05])

    add_metric_chip(slide, Inches(0.9), Inches(6.15), Inches(1.7), "Transactions", f"{counts.get('transactions', 0):,}", BLUE)
    add_metric_chip(slide, Inches(2.75), Inches(6.15), Inches(1.65), "Doctors", f"{counts.get('doctor_master', 0):,}", GREEN)
    add_metric_chip(slide, Inches(4.55), Inches(6.15), Inches(1.75), "Approvals", f"{counts.get('approval_requests', 0):,}", AMBER)
    add_metric_chip(slide, Inches(8.0), Inches(6.15), Inches(1.65), "Payments", f"{counts.get('payments', 0):,}", RED)
    add_metric_chip(slide, Inches(9.8), Inches(6.15), Inches(1.95), "Engine runs", f"{counts.get('engine_runs', 0):,}", BLUE)

    # Slide 11
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide)
    add_top_band(slide, "Monthly End-To-End Flow", "How the process starts and how it ends")
    add_page_title(slide, "Start-to-end monthly operating cycle", "This is the practical path the team follows each month from incoming files to finalized disbursal and reporting.")
    steps = [
        ("1. Setup", "Upload master files and verify doctor ownership."),
        ("2. Access", "Create users and link doctor logins where needed."),
        ("3. Intake", "Upload dashboard/incentive workbooks into transactions."),
        ("4. Validate", "Search, paginate, and export intake rows for checking."),
        ("5. Calculate", "Run RRCP engine to produce engine_runs and engine_results."),
        ("6. Review", "Check flags, variance, remarks, and productivity."),
        ("7. Payout", "Generate payments and capture finance adjustments."),
        ("8. Approve", "Approve doctor changes, overrides, and disbursal."),
        ("9. Close", "Disburse, lock the month, and preserve audit state."),
        ("10. Output", "Export reports or raise support corrections."),
    ]
    start_x = Inches(0.72)
    y = Inches(2.0)
    box_w = Inches(1.18)
    box_h = Inches(1.35)
    gap = Inches(0.09)
    for idx, (title, body) in enumerate(steps):
        x = start_x + idx * (box_w + gap)
        add_small_box(slide, x, y, box_w, box_h, title, body, fill=WHITE, accent=BLUE if idx < 4 else GREEN if idx < 7 else AMBER if idx < 9 else RED)
        if idx < len(steps) - 1:
            add_arrow(slide, x + box_w + Inches(0.01), y + Inches(0.46), Inches(0.08), Inches(0.36), "")

    add_body_card(
        slide,
        Inches(0.85),
        Inches(4.05),
        Inches(5.7),
        Inches(1.95),
        "Where the process starts",
        [
            "Reference masters and transaction workbooks are the primary starting inputs.",
            "Without updated doctor mapping and discount rules, the engine output is unreliable.",
        ],
        accent=GREEN,
    )
    add_body_card(
        slide,
        Inches(6.8),
        Inches(4.05),
        Inches(5.7),
        Inches(1.95),
        "Where the process ends",
        [
            "The month ends with approved payouts, locked period control, exported reports, and any unresolved issues parked in support.",
            "Every exception remains queryable through approvals, payment notes, and preserved raw intake rows.",
        ],
        accent=RED,
    )

    # Slide 12
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide)
    add_top_band(slide, "Next Improvements", "What is still worth doing after this delivery")
    add_page_title(slide, "Recommended next steps", "The app is operational, but a few targeted improvements would reduce future operational risk and manual cleanup.")
    add_body_card(
        slide,
        Inches(0.75),
        Inches(1.9),
        Inches(5.75),
        Inches(4.95),
        "Highest-value next work",
        [
            "Add duplicate-upload prevention so the same monthly workbook cannot inflate transactions twice.",
            "Add a dedicated audit log timeline for approval decisions, payout edits, and doctor profile changes.",
            "Add stricter upload validators with row-level error download for master files and monthly intake files.",
            "Add richer reports in Excel/PDF format beyond CSV if stakeholders need presentation-ready output.",
            "Add automated scheduled export or monthly deck generation once the workflow stabilizes.",
        ],
        accent=BLUE,
    )
    add_body_card(
        slide,
        Inches(6.8),
        Inches(1.9),
        Inches(5.75),
        Inches(4.95),
        "Status after this handoff",
        [
            "Workflow pages are compacted and aligned to the monthly process.",
            "Architecture screenshots are captured from the live local app state.",
            "This PPT can be regenerated using the scripts in scripts/capture_architecture_screenshots.py and scripts/build_architecture_ppt.py.",
            f"Generated output path: {OUTPUT_PATH.relative_to(ROOT)}",
        ],
        accent=GREEN,
    )

    ARTIFACT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    output = build_presentation()
    print(output)
